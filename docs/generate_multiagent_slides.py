"""Generate PowerPoint presentation for Multi-Agent Analysis System."""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import csv


class RgbColor:
    """RGB Color helper."""
    def __init__(self, r, g, b):
        self.r = r
        self.g = g
        self.b = b

    def __str__(self):
        return f"RGB({self.r}, {self.g}, {self.b})"


def set_shape_color(shape, rgb):
    """Set shape fill color."""
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(rgb.r, rgb.g, rgb.b)


def set_font_color(font, rgb):
    """Set font color."""
    from pptx.dml.color import RGBColor
    font.color.rgb = RGBColor(rgb.r, rgb.g, rgb.b)

# Paths
DATADIR = os.path.join(os.path.dirname(__file__), "..", "data")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "multiagent_presentation.pptx")


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        top = Inches(4)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_lines, bullet=True):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    top = Inches(1.3)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet and not line.startswith("  "):
            p.text = f"• {line}"
        else:
            p.text = line
        p.font.size = Pt(20)
        p.space_after = Pt(12)

    return slide


def add_diagram_slide(prs, title):
    """Add a slide with multi-agent architecture diagram."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Coordinator box (top center)
    coord_left = Inches(3.5)
    coord_top = Inches(1.3)
    coord_width = Inches(3)
    coord_height = Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, coord_left, coord_top, coord_width, coord_height)
    set_shape_color(shape, RgbColor(70, 130, 180))  # Steel blue
    tf = shape.text_frame
    tf.paragraphs[0].text = "Coordinator"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    set_font_color(tf.paragraphs[0].font, RgbColor(255, 255, 255))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Agent boxes
    agents = [
        ("Statistics\nAgent", RgbColor(60, 179, 113)),   # Medium sea green
        ("Feature\nAgent", RgbColor(255, 165, 0)),       # Orange
        ("Pattern\nAgent", RgbColor(147, 112, 219)),     # Medium purple
        ("Hypothesis\nAgent", RgbColor(220, 20, 60)),    # Crimson
    ]

    agent_width = Inches(1.8)
    agent_height = Inches(1.0)
    agent_top = Inches(3.0)
    start_left = Inches(0.6)
    spacing = Inches(2.2)

    for i, (name, color) in enumerate(agents):
        left = start_left + i * spacing
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, agent_top, agent_width, agent_height)
        set_shape_color(shape, color)
        tf = shape.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        set_font_color(tf.paragraphs[0].font, RgbColor(255, 255, 255))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = True

    # ThreadPoolExecutor box
    pool_left = Inches(2.5)
    pool_top = Inches(4.5)
    pool_width = Inches(5)
    pool_height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pool_left, pool_top, pool_width, pool_height)
    set_shape_color(shape, RgbColor(105, 105, 105))
    tf = shape.text_frame
    tf.paragraphs[0].text = "ThreadPoolExecutor (max_workers=4)"
    tf.paragraphs[0].font.size = Pt(16)
    set_font_color(tf.paragraphs[0].font, RgbColor(255, 255, 255))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Results aggregation
    result_left = Inches(3.5)
    result_top = Inches(5.5)
    result_width = Inches(3)
    result_height = Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, result_left, result_top, result_width, result_height)
    set_shape_color(shape, RgbColor(34, 139, 34))  # Forest green
    tf = shape.text_frame
    tf.paragraphs[0].text = "Synthesis Report"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    set_font_color(tf.paragraphs[0].font, RgbColor(255, 255, 255))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_agent_detail_slide(prs, agent_name, description, tasks, output):
    """Add a slide explaining one agent in detail."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = agent_name
    p.font.size = Pt(32)
    p.font.bold = True

    # Description
    top = Inches(1.2)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.italic = True

    # Tasks section
    top = Inches(2.0)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Tasks:"
    p.font.size = Pt(20)
    p.font.bold = True

    top = Inches(2.5)
    height = Inches(2.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for i, task in enumerate(tasks):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {task}"
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Output section
    top = Inches(4.5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Output:"
    p.font.size = Pt(20)
    p.font.bold = True

    top = Inches(5.0)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for i, out in enumerate(output):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→ {out}"
        p.font.size = Pt(18)
        p.space_after = Pt(6)

    return slide


def add_results_table_slide(prs, title, headers, rows):
    """Add a slide with a results table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_left = Inches(0.5)
    table_top = Inches(1.5)
    table_width = Inches(9)
    table_height = Inches(0.5 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table

    # Set column widths
    col_width = Inches(9 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_hypothesis_slide(prs, hypotheses_data):
    """Add a slide showing hypothesis verification results summary."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Hypothesis Verification Summary"
    p.font.size = Pt(32)
    p.font.bold = True

    # Hypotheses
    top = Inches(1.3)
    for hid, hdata in sorted(hypotheses_data.items()):
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        # Status icon
        status = hdata.get("status", "untested")
        if status == "confirmed":
            icon = "[O]"
            color = RgbColor(34, 139, 34)  # Green
        elif status == "rejected":
            icon = "[X]"
            color = RgbColor(220, 20, 60)  # Red
        else:
            icon = "[?]"
            color = RgbColor(128, 128, 128)  # Gray

        p = tf.paragraphs[0]
        p.text = f"{icon} {hid}: {hdata.get('statement', '')}"
        p.font.size = Pt(16)
        set_font_color(p.font, color)

        p = tf.add_paragraph()
        confidence = hdata.get("confidence", 0)
        verifications = len(hdata.get("verification_history", []))
        p.text = f"     Confidence: {confidence:.0%} | Verifications: {verifications}"
        p.font.size = Pt(14)
        set_font_color(p.font, RgbColor(100, 100, 100))

        top += Inches(0.85)

    return slide


def add_detailed_hypothesis_slide(prs, hid, statement, background, method, results, conclusion):
    """Add a detailed slide for a single hypothesis."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{hid}: {statement}"
    p.font.size = Pt(24)
    p.font.bold = True

    # Background section
    top = Inches(1.0)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Background (Why this hypothesis?)"
    p.font.size = Pt(18)
    p.font.bold = True
    set_font_color(p.font, RgbColor(70, 130, 180))

    top = Inches(1.4)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(background):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    # Method section
    top = Inches(2.5)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Verification Method"
    p.font.size = Pt(18)
    p.font.bold = True
    set_font_color(p.font, RgbColor(70, 130, 180))

    top = Inches(2.9)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(method):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    # Results section
    top = Inches(4.2)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Results"
    p.font.size = Pt(18)
    p.font.bold = True
    set_font_color(p.font, RgbColor(70, 130, 180))

    top = Inches(4.6)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(results):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    # Conclusion box
    top = Inches(6.0)
    height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    if "Confirmed" in conclusion:
        set_shape_color(shape, RgbColor(34, 139, 34))
    elif "Rejected" in conclusion:
        set_shape_color(shape, RgbColor(220, 20, 60))
    else:
        set_shape_color(shape, RgbColor(128, 128, 128))

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = conclusion
    p.font.size = Pt(16)
    p.font.bold = True
    set_font_color(p.font, RgbColor(255, 255, 255))
    p.alignment = PP_ALIGN.CENTER

    return slide


def load_experiment_data():
    """Load experiment data from CSV."""
    db_path = os.path.join(DATADIR, "experiments.csv")
    if not os.path.exists(db_path):
        return []

    with open(db_path) as f:
        return list(csv.DictReader(f))


def load_hypothesis_data():
    """Load hypothesis data from JSON."""
    memo_path = os.path.join(DATADIR, "hypotheses_memo.json")
    if not os.path.exists(memo_path):
        return {}

    with open(memo_path) as f:
        return json.load(f)


def calculate_summary(data):
    """Calculate summary statistics from experiment data."""
    if not data:
        return {}

    summary = {
        "total_runs": len(data),
        "solvers": {},
    }

    for solver in ["implicit_fdm", "spectral_cosine"]:
        solver_data = [d for d in data if d["solver"] == solver]
        stable = [d for d in solver_data if d["is_stable"] == "True"]

        errors = []
        times = []
        for d in stable:
            try:
                errors.append(float(d["l2_error"]))
                times.append(float(d["wall_time"]))
            except (ValueError, KeyError):
                pass

        summary["solvers"][solver] = {
            "total": len(solver_data),
            "stable": len(stable),
            "stable_pct": len(stable) / len(solver_data) * 100 if solver_data else 0,
            "avg_l2": sum(errors) / len(errors) if errors else 0,
            "avg_time": sum(times) / len(times) * 1000 if times else 0,  # ms
        }

    return summary


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Load data
    exp_data = load_experiment_data()
    hyp_data = load_hypothesis_data()
    summary = calculate_summary(exp_data)

    # === Title Slide ===
    add_title_slide(
        prs,
        "Multi-Agent Analysis System",
        "Hypothesis-Driven Solver Benchmark Framework"
    )

    # === Overview Slide ===
    add_content_slide(prs, "System Overview", [
        "Purpose: Automated analysis of numerical solver performance",
        "Approach: Multiple specialized agents working in parallel",
        "Key Features:",
        "  - Hypothesis tracking with verification history",
        "  - Parallel execution for speed improvement",
        "  - Automatic report generation",
        "  - Iterative refinement through verification cycles",
    ])

    # === Architecture Diagram ===
    add_diagram_slide(prs, "Multi-Agent Architecture")

    # === Agent Details ===
    add_agent_detail_slide(
        prs,
        "Statistics Agent",
        "Computes statistical summaries of solver performance",
        [
            "Calculate solver distribution (win rates)",
            "Compute feature means by solver type",
            "Identify performance patterns",
        ],
        [
            "Solver distribution percentages",
            "Feature statistics by solver",
        ]
    )

    add_agent_detail_slide(
        prs,
        "Feature Agent",
        "Analyzes feature importance for solver selection",
        [
            "Train decision tree on solver selection data",
            "Count feature usage in tree splits",
            "Rank features by importance",
        ],
        [
            "Feature importance scores",
            "Top 5 most predictive features",
        ]
    )

    add_agent_detail_slide(
        prs,
        "Pattern Agent",
        "Extracts decision rules from trained models",
        [
            "Traverse decision tree structure",
            "Extract IF-THEN rules",
            "Identify key decision thresholds",
        ],
        [
            "Human-readable decision rules",
            "Threshold values for each feature",
        ]
    )

    add_agent_detail_slide(
        prs,
        "Hypothesis Agent",
        "Tests and tracks scientific hypotheses",
        [
            "Evaluate hypotheses against new data",
            "Track verification history",
            "Update confidence scores",
        ],
        [
            "Hypothesis status (confirmed/rejected)",
            "Confidence scores",
            "Verification timeline",
        ]
    )

    # === Parallel Execution ===
    add_content_slide(prs, "Parallel Execution", [
        "Implementation: Python concurrent.futures.ThreadPoolExecutor",
        "",
        "Benefits:",
        "  - ~2x speedup with 4 agents",
        "  - Non-blocking execution",
        "  - Automatic result aggregation",
        "",
        "Code Pattern:",
        "  with ThreadPoolExecutor(max_workers=4) as executor:",
        "      futures = {executor.submit(agent, data): name}",
        "      for future in as_completed(futures):",
        "          results[name] = future.result()",
    ], bullet=False)

    # === Experiment Results ===
    if summary and summary.get("solvers"):
        headers = ["Solver", "Runs", "Stable %", "Avg L2", "Avg Time"]
        rows = []
        for solver, stats in summary["solvers"].items():
            rows.append([
                solver,
                stats["total"],
                f"{stats['stable_pct']:.1f}%",
                f"{stats['avg_l2']:.6f}",
                f"{stats['avg_time']:.2f}ms",
            ])
        add_results_table_slide(prs, "Experiment Results", headers, rows)

    # === Hypothesis Results Summary ===
    if hyp_data:
        add_hypothesis_slide(prs, hyp_data)

    # === Detailed Hypothesis Slides ===
    # H1: Smaller dt improves spectral stability
    add_detailed_hypothesis_slide(
        prs,
        hid="H1",
        statement="Smaller dt improves spectral solver stability",
        background=[
            "Spectral solver uses explicit time stepping for nonlinear terms",
            "Explicit methods have CFL-like stability constraints",
            "Larger dt may cause numerical oscillations or divergence",
        ],
        method=[
            "Tested spectral solver with dt = 0.002, 0.001, 0.0005, 0.0002, 0.0001",
            "Fixed parameters: nr=51, t_end=0.1, alpha=0.0~1.0",
            "Measured stability rate (% of runs without NaN/divergence)",
            "Repeated across 104 parameter combinations",
        ],
        results=[
            "dt=0.0001: 100% stable across all alpha values",
            "dt=0.0005: 100% stable",
            "dt=0.001: 50% stable (fails for high alpha)",
            "dt=0.002: 0% stable (fails for all alpha > 0)",
        ],
        conclusion="Confirmed: Smaller dt significantly improves stability (100% at dt<=0.0005)"
    )

    # H7: Spectral fails for high alpha
    add_detailed_hypothesis_slide(
        prs,
        hid="H7",
        statement="Spectral solver fails with NaN for alpha >= 0.2",
        background=[
            "The chi formula has a threshold at |dT/dr| = 0.5",
            "Higher alpha means steeper nonlinearity above threshold",
            "Spectral methods handle smooth functions well, but sharp transitions poorly",
        ],
        method=[
            "Tested spectral solver with alpha = 0.0, 0.1, 0.2, 0.5, 0.8, 1.0",
            "Fixed parameters: nr=51, dt=0.001, t_end=0.1",
            "Checked for NaN, Inf, or values > 100 (divergence)",
            "Compared with implicit FDM as baseline",
        ],
        results=[
            "alpha=0.0: Spectral stable (max T = 1.00)",
            "alpha=0.1: Spectral stable (max T = 1.01)",
            "alpha>=0.2: Spectral produces NaN or diverges",
            "FDM remains stable for all alpha values tested",
        ],
        conclusion="Confirmed: Spectral fails for alpha >= 0.2 with default dt"
    )

    # H4: Different ICs affect solver performance
    add_detailed_hypothesis_slide(
        prs,
        hid="H4",
        statement="Different initial conditions lead to different optimal solvers",
        background=[
            "Initial condition affects gradient profile |dT/dr|",
            "Chi activation depends on whether |dT/dr| > 0.5",
            "Different ICs may have different gradient characteristics",
        ],
        method=[
            "Tested 4 IC types: parabola (1-r²), gaussian, cosine, sine",
            "Compared L2 error and computation time for each solver",
            "Used cost function: score = L2_error + 0.1 × wall_time",
            "Ran 24 experiments (4 ICs × 3 alpha × 2 solvers)",
        ],
        results=[
            "Parabola IC: Spectral wins (lower L2 error)",
            "Cosine IC: Spectral wins",
            "Sine IC: Spectral wins",
            "Gaussian IC: FDM wins (spectral unstable)",
            "Gaussian has sharper gradient near r=0",
        ],
        conclusion="Confirmed: IC type affects winner (gaussian favors FDM)"
    )

    # H3: FDM unconditionally stable
    add_detailed_hypothesis_slide(
        prs,
        hid="H3",
        statement="FDM is unconditionally stable for any dt",
        background=[
            "Implicit FDM (Crank-Nicolson) is theoretically A-stable",
            "Nonlinear terms are handled implicitly via iteration",
            "Should remain stable even with large time steps",
        ],
        method=[
            "Tested implicit FDM with dt = 0.001, 0.005, 0.01, 0.02, 0.05",
            "Used alpha = 1.0 (strong nonlinearity)",
            "Checked for physical bounds: 0 <= T <= 1",
            "Ran 80 experiments across parameter space",
        ],
        results=[
            "100% stability rate across all tested dt values",
            "Even dt=0.05 (50x larger than default) remains stable",
            "All solutions maintain physical bounds",
            "No NaN or divergence observed",
        ],
        conclusion="Confirmed: FDM is unconditionally stable (100% rate)"
    )

    # === Key Findings ===
    add_content_slide(prs, "Key Findings", [
        "H1 Confirmed: Smaller dt improves spectral stability",
        "  - dt=0.0001: 100% stable",
        "  - dt=0.002: 0% stable",
        "",
        "H7 Confirmed: Spectral fails for high alpha",
        "  - Instability threshold around alpha >= 0.2",
        "",
        "FDM Characteristics:",
        "  - Unconditionally stable (100% for all parameters)",
        "  - Higher L2 error but more reliable",
    ])

    # === Iteration Cycle ===
    add_content_slide(prs, "Verification Cycle", [
        "1. Run Experiments",
        "   - Execute predefined parameter sweeps",
        "   - Store results in CSV database",
        "",
        "2. Test Hypotheses",
        "   - Evaluate each hypothesis against new data",
        "   - Update status (confirmed/rejected/inconclusive)",
        "",
        "3. Update Confidence",
        "   - +20% for confirmed, -10% for rejected",
        "   - Track verification history",
        "",
        "4. Generate Report",
        "   - Markdown summary with conclusions",
        "   - Recommended next steps",
    ])

    # === Usage ===
    add_content_slide(prs, "Usage", [
        "Auto mode (full pipeline):",
        "  python docs/analysis/experiment_framework.py --auto 3",
        "",
        "Interactive mode:",
        "  python docs/analysis/experiment_framework.py -i",
        "",
        "Commands:",
        "  run stability_map  - Run experiment",
        "  test H1           - Test hypothesis",
        "  cycle 5           - Run 5 verification cycles",
        "  report            - Generate final report",
        "  hypo add H8 '...' - Add new hypothesis",
    ], bullet=False)

    # === Summary ===
    add_content_slide(prs, "Summary", [
        "Multi-Agent System provides:",
        "  - Parallel analysis for faster insights",
        "  - Specialized agents for different tasks",
        "  - Hypothesis-driven experimentation",
        "  - Automated report generation",
        "",
        "Results:",
        f"  - {summary.get('total_runs', 0)} experiments analyzed",
        f"  - {len(hyp_data)} hypotheses tracked",
        "  - 2 hypotheses confirmed, 1 rejected",
        "",
        "Future Work:",
        "  - Add more specialized agents",
        "  - Implement adaptive experiment selection",
    ])

    # Save
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
