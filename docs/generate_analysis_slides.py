"""Generate analysis slides for feature-cost function relationship."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGDIR = os.path.join(os.path.dirname(__file__), "figures")
OUTPATH = os.path.join(os.path.dirname(__file__), "solver_analysis.pptx")

# Color palette
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_MID = RGBColor(0x24, 0x24, 0x3E)
ACCENT = RGBColor(0x4E, 0xA8, 0xDE)
ACCENT2 = RGBColor(0x50, 0xC8, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
RED = RGBColor(0xE0, 0x50, 0x50)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_code_box(slide, left, top, width, height, code, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x28)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xA0, 0xE0, 0xA0)
    p.font.name = "Consolas"
    return tf


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
                 "Solver Selection Analysis", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.7),
                 "Feature-Cost Function Relationship", font_size=28, color=ACCENT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(8.4), Inches(1.0),
                 "Analysis of which solver wins under which conditions\n"
                 "based on 432 parameter combinations",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.5),
                 "score = L2_error + 0.1 \u00d7 wall_time",
                 font_size=22, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")


def slide_cost_function(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Cost Function Definition", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(0.6),
                 "score = L2_error + \u03bb \u00d7 wall_time",
                 font_size=28, color=ORANGE, font_name="Cambria Math",
                 alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(3.5), [
        "\u25b6  L2_error: Weighted L2 norm vs reference solution",
        "      L2 = \u221a(\u222b(T - T_ref)\u00b2 r dr / \u222b T_ref\u00b2 r dr)",
        "",
        "\u25b6  wall_time: Actual computation time [seconds]",
        "",
        "\u25b6  \u03bb = 0.1 (default): Accuracy-focused selection",
        "      \u03bb = 0 \u2192 pure accuracy, \u03bb = 1 \u2192 speed-focused",
        "",
        "\u25b6  Reference: Same Implicit FDM with 4\u00d7 refinement",
    ], font_size=17)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.5),
                 "The solver with the LOWEST score wins",
                 font_size=20, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "12 Features for Solver Prediction", font_size=32, color=ACCENT, bold=True)

    # Three columns for categories
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(3.0), Inches(0.4),
                 "Problem Params (4)", font_size=16, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(1.6), Inches(3.0), Inches(2.0), [
        "\u03b1 (alpha)",
        "nr (grid points)",
        "dt (time step)",
        "t_end (sim time)",
    ], font_size=14)

    add_text_box(slide, Inches(3.5), Inches(1.2), Inches(3.0), Inches(0.4),
                 "Physical Features (5)", font_size=16, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(3.5), Inches(1.6), Inches(3.0), Inches(2.5), [
        "max_abs_gradient",
        "energy_content",
        "max_chi",
        "max_laplacian",
        "T_center",
    ], font_size=14)

    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(3.3), Inches(0.4),
                 "Derived Features (3)", font_size=16, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(6.5), Inches(1.6), Inches(3.3), Inches(2.0), [
        "gradient_sharpness",
        "chi_ratio",
        "problem_stiffness",
    ], font_size=14)

    # Key insight
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(0.5),
                 "Key Derived Feature:", font_size=18, color=ORANGE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.8),
                 "problem_stiffness = \u03b1 \u00d7 max_abs_gradient\n"
                 "Captures interaction between nonlinearity and gradient steepness",
                 font_size=16, color=WHITE)

    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9.0), Inches(0.5),
                 "For T\u2080 = 1 - r\u00b2:", font_size=18, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(5.9), Inches(8.4), Inches(1.2), [
        "max_abs_gradient \u2248 2.0 (at r=1)",
        "max_chi = (\u03b1)(1.5)\u03b1 + 0.1 (since |T'| = 2r > 0.5 for r > 0.25)",
    ], font_size=15)


def slide_win_distribution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Analysis Result: Solver Win Distribution", font_size=32, color=ACCENT, bold=True)

    fig = os.path.join(FIGDIR, "solver_distribution.png")
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(0.5), Inches(1.1), Inches(5.5), Inches(3.5))

    # Key findings
    add_text_box(slide, Inches(6.2), Inches(1.2), Inches(3.5), Inches(0.4),
                 "Key Findings", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(6.2), Inches(1.7), Inches(3.6), Inches(3.0), [
        "implicit_fdm: 99.5%",
        "  (430 / 432 cases)",
        "",
        "spectral_cosine: 0.5%",
        "  (2 / 432 cases)",
        "",
        "pinn_stub: 0%",
        "  (NaN errors)",
    ], font_size=15)

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9.0), Inches(1.5),
                 "Conclusion: Implicit FDM dominates due to unconditional stability\n"
                 "of Crank-Nicolson scheme with the threshold-based \u03c7 formula.",
                 font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_feature_importance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Feature Importance in Decision Tree", font_size=32, color=ACCENT, bold=True)

    fig = os.path.join(FIGDIR, "feature_importance.png")
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(0.3), Inches(1.0), Inches(6.0), Inches(3.8))

    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(3.3), Inches(0.4),
                 "Top 5 Features", font_size=18, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(6.5), Inches(1.7), Inches(3.3), Inches(2.5), [
        "1. dt (0.29)",
        "2. nr (0.14)",
        "3. t_end (0.14)",
        "4. max_abs_gradient (0.14)",
        "5. max_chi (0.14)",
    ], font_size=15)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(1.5),
                 "Interpretation: Time step (dt) is the most critical feature.\n"
                 "The spectral solver requires finer time steps for stability.",
                 font_size=17, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_alpha_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Solver Performance vs Alpha", font_size=32, color=ACCENT, bold=True)

    fig = os.path.join(FIGDIR, "alpha_analysis.png")
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(0.3), Inches(1.0), Inches(9.4), Inches(3.3))

    add_bullet_list(slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(2.5), [
        "\u25b6  Spectral solver only wins at \u03b1 = 0.0 and \u03b1 = 0.1",
        "\u25b6  For \u03b1 \u2265 0.2, implicit FDM wins 100% of cases",
        "\u25b6  Higher \u03b1 \u2192 stronger nonlinearity \u2192 spectral instability",
        "\u25b6  The threshold \u03c7 formula causes numerical issues in spectral method",
    ], font_size=17)


def slide_spectral_conditions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "When Does Spectral Solver Win?", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.5),
                 "Spectral wins in only 2 out of 432 cases:", font_size=20, color=ACCENT2, bold=True)

    # Table-like layout
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.5),
                 "Case 1:  \u03b1=0.0, nr=71, dt=0.001, t_end=0.05\n"
                 "Case 2:  \u03b1=0.1, nr=31, dt=0.002, t_end=0.1\n\n"
                 "Common characteristics:\n"
                 "  \u2022 Very low alpha (\u03b1 \u2264 0.1) \u2192 near-linear diffusion\n"
                 "  \u2022 problem_stiffness < 0.27 \u2192 weak nonlinearity\n"
                 "  \u2022 chi_ratio \u2248 11 \u2192 minimal diffusivity variation",
                 font_size=17, color=WHITE)

    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.5),
                 "Decision Rule from Tree:", font_size=20, color=ORANGE, bold=True)
    add_code_box(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(1.5),
                 "IF problem_stiffness <= 0.33 AND max_abs_gradient <= 2.71:\n"
                 "    THEN spectral_cosine MAY win\n"
                 "ELSE:\n"
                 "    THEN implicit_fdm wins",
                 font_size=14)


def slide_dt_nr_heatmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "FDM Win Rate by dt and nr", font_size=32, color=ACCENT, bold=True)

    fig = os.path.join(FIGDIR, "dt_nr_heatmap.png")
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(1.0), Inches(1.0), Inches(6.0), Inches(4.5))

    add_text_box(slide, Inches(7.2), Inches(1.5), Inches(2.5), Inches(3.0),
                 "Observation:\n\n"
                 "100% FDM wins\n"
                 "across all\n"
                 "dt-nr combinations\n\n"
                 "(Blue = FDM wins)",
                 font_size=15, color=LIGHT_GRAY)


def slide_why_fdm_wins(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Why Does Implicit FDM Dominate?", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.4),
                 "1. Numerical Stability", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(1.6), Inches(4.5), Inches(1.5), [
        "Crank-Nicolson is unconditionally stable",
        "Handles stiff nonlinear \u03c7 correctly",
        "L'H\u00f4pital rule at r=0 is clean",
    ], font_size=15)

    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(0.4),
                 "2. Spectral Solver Issues", font_size=20, color=RED, bold=True)
    add_bullet_list(slide, Inches(5.2), Inches(1.6), Inches(4.5), Inches(1.5), [
        "Operator splitting is semi-explicit",
        "Threshold \u03c7 causes Gibbs-like issues",
        "Unstable for \u03b1 > 0.1",
    ], font_size=15)

    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(9.0), Inches(0.4),
                 "3. Cost Function Favors Accuracy", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(3.7), Inches(9.0), Inches(1.5), [
        "\u03bb = 0.1 means error matters 10\u00d7 more than time",
        "FDM is slightly slower but much more accurate",
        "Spectral's NaN/Inf errors \u2192 automatic disqualification",
    ], font_size=15)

    add_text_box(slide, Inches(0.5), Inches(5.3), Inches(9.0), Inches(0.4),
                 "4. Initial Condition T\u2080 = 1 - r\u00b2", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(5.7), Inches(9.0), Inches(1.0), [
        "max|dT/dr| = 2 at r=1 \u2192 always above threshold 0.5",
        "Nonlinear regime is active for all cases",
    ], font_size=15)


def slide_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Recommendations", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.4),
                 "For Users:", font_size=22, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(1.7), Inches(9.0), Inches(1.5), [
        "\u2705  Use implicit_fdm as the default solver",
        "\u2705  Consider spectral only for \u03b1 < 0.1 (near-linear regime)",
        "\u2705  ML selector is not needed for current IC (T\u2080 = 1 - r\u00b2)",
    ], font_size=17)

    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9.0), Inches(0.4),
                 "For Development:", font_size=22, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(2.0), [
        "\u26a0  Improve spectral solver stability for threshold \u03c7",
        "\u26a0  Consider implicit treatment of nonlinear term in spectral",
        "\u26a0  Add more IC types to create diverse training data",
        "\u26a0  Implement full PINN to provide a third viable option",
    ], font_size=17)

    add_text_box(slide, Inches(0.5), Inches(6.0), Inches(9.0), Inches(0.5),
                 "The ML selector will become useful when solvers have\n"
                 "more balanced performance across parameter space.",
                 font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Summary", font_size=32, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(4.5), [
        "\u25b6  Analyzed 432 parameter combinations (8\u03b1 \u00d7 3nr \u00d7 3dt \u00d7 3t_end)",
        "",
        "\u25b6  Cost function: score = L2_error + 0.1 \u00d7 wall_time",
        "",
        "\u25b6  Implicit FDM wins 99.5% of cases due to:",
        "      \u2022 Unconditional stability of Crank-Nicolson",
        "      \u2022 Robust handling of threshold-based \u03c7",
        "",
        "\u25b6  Spectral solver only competitive at \u03b1 \u2264 0.1",
        "",
        "\u25b6  Most important feature: dt (time step)",
        "",
        "\u25b6  Decision rule: IF problem_stiffness > 0.33 THEN use FDM",
    ], font_size=17)


def slide_tools(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Analysis Tools Created", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.4),
                 "New Script: docs/analyze_solver_selection.py", font_size=18, color=ACCENT2, bold=True)

    add_bullet_list(slide, Inches(0.5), Inches(1.7), Inches(9.0), Inches(2.5), [
        "\u25b6  Generates/loads training data (432 samples)",
        "\u25b6  Computes solver win distribution",
        "\u25b6  Analyzes feature statistics by winning solver",
        "\u25b6  Trains decision tree and extracts feature importance",
        "\u25b6  Extracts human-readable decision rules",
        "\u25b6  Generates 6 visualization figures",
    ], font_size=16)

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(9.0), Inches(0.4),
                 "Generated Figures:", font_size=18, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(4.7), Inches(4.5), Inches(2.0), [
        "solver_distribution.png",
        "feature_importance.png",
        "alpha_analysis.png",
    ], font_size=14)
    add_bullet_list(slide, Inches(5.0), Inches(4.7), Inches(4.5), Inches(2.0), [
        "dt_nr_heatmap.png",
        "cost_breakdown.png",
        "alpha_vs_stiffness.png",
    ], font_size=14)

    add_code_box(slide, Inches(0.5), Inches(6.2), Inches(9.0), Inches(0.8),
                 "# Run analysis\npython docs/analyze_solver_selection.py",
                 font_size=13)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_cost_function(prs)
    slide_features(prs)
    slide_win_distribution(prs)
    slide_feature_importance(prs)
    slide_alpha_analysis(prs)
    slide_spectral_conditions(prs)
    slide_dt_nr_heatmap(prs)
    slide_why_fdm_wins(prs)
    slide_recommendations(prs)
    slide_summary(prs)
    slide_tools(prs)

    prs.save(OUTPATH)
    print(f"Saved {len(prs.slides)} slides to {OUTPATH}")


if __name__ == "__main__":
    main()
