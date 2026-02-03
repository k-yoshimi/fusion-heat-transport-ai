"""Generate presentation slides (PPTX) for the Fusion Heat Transport project."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGDIR = os.path.join(os.path.dirname(__file__), "figures")
OUTPATH = os.path.join(os.path.dirname(__file__), "presentation.pptx")

# Color palette
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_MID = RGBColor(0x24, 0x24, 0x3E)
ACCENT = RGBColor(0x4E, 0xA8, 0xDE)
ACCENT2 = RGBColor(0x50, 0xC8, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_code_box(slide, left, top, width, height, code, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x28)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xA0, 0xE0, 0xA0)
    p.font.name = "Consolas"
    return tf


def slide_title_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.0),
                 "Fusion Heat Transport", font_size=40, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8.4), Inches(0.7),
                 "PDE Benchmark", font_size=36, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.8),
                 "Numerical Solver Comparison with ML-Based Selection\n"
                 "for 1D Radial Heat Equation in Fusion Plasmas",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Equation
    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(7.0), Inches(0.5),
                 "\u2202T/\u2202t = (1/r) \u2202/\u2202r (r \u03c7(|\u2202T/\u2202r|) \u2202T/\u2202r)",
                 font_size=22, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")
    add_text_box(slide, Inches(1.5), Inches(4.9), Inches(7.0), Inches(0.5),
                 "\u03c7(|T'|) = (|T'| \u2212 0.5)^\u03b1 + 0.1  (|T'| > 0.5)",
                 font_size=18, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Project Overview", font_size=32, color=ACCENT, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(4.5), [
        "\u25b6  1D radial heat equation with nonlinear diffusivity \u03c7(|T'|)",
        "\u25b6  Models anomalous heat transport in fusion plasma devices",
        "\u25b6  8 numerical solvers: FDM, FEM, FVM, Spectral, PINN",
        "\u25b6  Automated benchmarking with L2/L\u221e error metrics",
        "\u25b6  ML-based solver selector (decision tree) predicts best solver",
        "\u25b6  Optimized implementations (up to 10x speedup)",
        "\u25b6  Pure Python — numpy + scipy only (no sklearn dependency)",
    ], font_size=18)


def slide_pde(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "The PDE: Radial Heat Equation", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
                 "\u2202T/\u2202t = (1/r) \u2202/\u2202r (r \u03c7 \u2202T/\u2202r)",
                 font_size=22, color=ORANGE, font_name="Cambria Math",
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(0.7),
                 "\u03c7(|T'|) = (|T'| \u2212 0.5)^\u03b1 + 0.1  if |T'| > 0.5\n"
                 "\u03c7(|T'|) = 0.1                    if |T'| \u2264 0.5",
                 font_size=17, color=ORANGE, font_name="Cambria Math",
                 alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(4.0), Inches(2.5), [
        "Boundary conditions:",
        "  r = 0:  \u2202T/\u2202r = 0  (symmetry)",
        "  r = 1:  T = 0  (fixed wall)",
        "",
        "Singularity at r = 0:",
        "  L'H\u00f4pital \u2192 2\u03c7 \u2202\u00b2T/\u2202r\u00b2",
    ], font_size=16)

    fig = os.path.join(FIGDIR, "nonlinear_diffusivity.png")
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(5.0), Inches(1.9), Inches(4.8), Inches(2.7))


def slide_initial_conditions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Initial Conditions & Time Evolution", font_size=32, color=ACCENT, bold=True)

    fig1 = os.path.join(FIGDIR, "initial_conditions.png")
    fig2 = os.path.join(FIGDIR, "time_evolution.png")
    if os.path.exists(fig1):
        slide.shapes.add_picture(fig1, Inches(0.3), Inches(1.2), Inches(4.5), Inches(2.8))
    if os.path.exists(fig2):
        slide.shapes.add_picture(fig2, Inches(0.3), Inches(4.1), Inches(9.4), Inches(2.8))

    add_bullet_list(slide, Inches(5.0), Inches(1.3), Inches(4.7), Inches(2.5), [
        "T\u2080(r) = 1 \u2212 r\u00b2 (parabolic)",
        "Satisfies both BCs: T\u2080(1)=0, dT\u2080/dr|_{r=0}=0",
        "\u03b1 = 0: standard linear diffusion",
        "\u03b1 > 0: enhanced diffusion at steep gradients",
    ], font_size=14)


def slide_solvers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Numerical Solvers (8 Methods)", font_size=32, color=ACCENT, bold=True)

    # Table-like layout with performance
    headers = ["Solver", "Method", "Time", "L2 Err"]
    rows = [
        ["cell_centered_fvm", "Finite Volume", "1.9ms", "0.058"],
        ["implicit_fdm", "Crank-Nicolson FDM", "2.2ms", "0.083"],
        ["compact4_fdm", "4th-order Compact FDM", "2.5ms", "0.115"],
        ["imex_fdm", "IMEX (split operator)", "3.4ms", "0.492"],
        ["chebyshev_spectral", "Chebyshev Spectral", "3.6ms", "0.799"],
        ["p2_fem", "P2 Finite Element", "39.4ms", "0.108"],
        ["cosine_spectral", "Cosine Expansion", "2.5ms", "varies"],
        ["pinn_stub", "Physics-Informed NN", "-", "-"],
    ]

    y = Inches(1.0)
    col_specs = [(Inches(0.4), Inches(2.4)), (Inches(2.8), Inches(2.8)),
                 (Inches(5.6), Inches(1.2)), (Inches(6.9), Inches(1.2))]
    for col, (x, w) in enumerate(col_specs):
        add_text_box(slide, x, y, w, Inches(0.35),
                     headers[col], font_size=14, color=ACCENT, bold=True)

    for i, row in enumerate(rows):
        y_row = Inches(1.4 + i * 0.42)
        for col, (x, w) in enumerate(col_specs):
            c = ORANGE if col == 0 else (ACCENT2 if col == 2 else WHITE)
            add_text_box(slide, x, y_row, w, Inches(0.38),
                         row[col], font_size=12, color=c)

    # Reference solution
    add_text_box(slide, Inches(0.4), Inches(5.0), Inches(9.0), Inches(0.4),
                 "Reference: Implicit FDM with 4\u00d7 refinement (nr\u00d74, dt/4)",
                 font_size=14, color=LIGHT_GRAY)


def slide_benchmark_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Benchmark Results", font_size=32, color=ACCENT, bold=True)

    fig1 = os.path.join(FIGDIR, "solver_comparison.png")
    fig2 = os.path.join(FIGDIR, "alpha_sweep.png")
    if os.path.exists(fig1):
        slide.shapes.add_picture(fig1, Inches(0.2), Inches(1.1), Inches(9.6), Inches(3.0))
    if os.path.exists(fig2):
        slide.shapes.add_picture(fig2, Inches(0.2), Inches(4.2), Inches(9.6), Inches(3.0))


def slide_selection_policy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Solver Selection Policy", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(0.5),
                 "score = L2_error + \u03bb \u00d7 wall_time",
                 font_size=24, color=ORANGE, font_name="Cambria Math",
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(2.2), Inches(4.5), Inches(0.4),
                 "Post-hoc Selection (default)", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(4.3), Inches(2.0), [
        "Run ALL solvers",
        "Compute errors vs reference",
        "Pick lowest score",
        "\u03bb = 0.1 (accuracy-focused)",
    ], font_size=16)

    add_text_box(slide, Inches(5.2), Inches(2.2), Inches(4.5), Inches(0.4),
                 "ML Prediction (new)", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(5.4), Inches(2.7), Inches(4.3), Inches(2.0), [
        "Extract 14 features from T\u2080",
        "Decision tree predicts best solver",
        "Run ONLY predicted solver",
        "No sklearn \u2014 numpy-only tree",
    ], font_size=16)


def slide_ml_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "ML Feature Engineering (12 Features)", font_size=32, color=ACCENT, bold=True)

    categories = [
        ("Problem Parameters (4)", ACCENT2, [
            "\u03b1, nr, dt, t_end",
        ]),
        ("Physical Features from T\u2080 (5)", ACCENT2, [
            "max_abs_gradient, energy_content",
            "max_chi, max_laplacian, T_center",
        ]),
        ("Derived Features (3)", ACCENT2, [
            "gradient_sharpness = max_grad / T_center",
            "chi_ratio = max_chi / min_chi",
            "problem_stiffness = \u03b1 \u00d7 max_grad",
        ]),
    ]

    y = Inches(1.2)
    for title, color, items in categories:
        add_text_box(slide, Inches(0.8), y, Inches(8.4), Inches(0.4),
                     title, font_size=18, color=color, bold=True)
        y += Inches(0.45)
        for item in items:
            add_text_box(slide, Inches(1.2), y, Inches(8.0), Inches(0.35),
                         "\u2022  " + item, font_size=15, color=WHITE)
            y += Inches(0.35)
        y += Inches(0.15)

    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(8.4), Inches(0.5),
                 "All features extracted BEFORE solving \u2192 enables prediction without computation",
                 font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_ml_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "ML Workflow", font_size=32, color=ACCENT, bold=True)

    # 3-step flow
    steps = [
        ("1", "Generate Data", "Parameter sweep\n~216 instances\n8 \u03b1 \u00d7 3 nr\n\u00d7 3 dt \u00d7 3 t_end"),
        ("2", "Train Model", "CART decision tree\nGini impurity\nmax_depth = 5\nnumpy-only"),
        ("3", "Predict & Run", "Extract features\nPredict best solver\nRun only that one"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 3.3)
        # Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Inches(1.5), Inches(2.8), Inches(2.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_MID
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(2.4), Inches(0.5),
                     f"Step {num}", font_size=14, color=LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.4), Inches(0.5),
                     title, font_size=20, color=ACCENT2, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.6), Inches(2.4), Inches(1.5),
                     desc, font_size=14, color=WHITE)

    # Arrow symbols between boxes
    for i in range(2):
        x = Inches(3.3 + i * 3.3)
        add_text_box(slide, x, Inches(2.5), Inches(0.5), Inches(0.5),
                     "\u25b6", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # CLI commands
    add_code_box(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(2.0),
                 "# Step 1: Generate training data\n"
                 "python -m app.run_benchmark --generate-data\n\n"
                 "# Step 2: Train model\n"
                 "python -m policy.train    # or: make train\n\n"
                 "# Step 3: Use ML selector\n"
                 "python -m app.run_benchmark --use-ml-selector --alpha 1.5",
                 font_size=13)


def slide_incremental(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Incremental Learning with --update", font_size=32, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(2.5), [
        "\u25b6  Each benchmark run can improve the ML model",
        "\u25b6  --update flag: run benchmark \u2192 append results to CSV \u2192 retrain",
        "\u25b6  Model accumulates experience over time",
        "\u25b6  No need to regenerate full training set",
    ], font_size=18)

    add_code_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(1.8),
                 "# Run benchmark and update model\n"
                 "python -m app.run_benchmark --alpha 0.3 --init sharp --update\n\n"
                 "# Accumulate more data\n"
                 "python -m app.run_benchmark --alpha 1.2 --nr 71 --update\n\n"
                 "# Use improved model\n"
                 "python -m app.run_benchmark --use-ml-selector --alpha 0.3",
                 font_size=13)

    # Flow diagram text
    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(8.4), Inches(0.5),
                 "Benchmark  \u2192  Label best  \u2192  Append CSV  \u2192  Retrain tree  \u2192  Updated model",
                 font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_optimization(prs):
    """Slide showing optimization results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Performance Optimization", font_size=32, color=ACCENT, bold=True)

    # Optimization table
    headers = ["Solver", "Before", "After", "Speedup"]
    rows = [
        ["Compact4 FDM", "25.3ms", "2.5ms", "10.3\u00d7"],
        ["Cell-Centered FVM", "14.5ms", "1.9ms", "7.6\u00d7"],
        ["P2 FEM", "285.0ms", "39.4ms", "7.2\u00d7"],
        ["Chebyshev Spectral", "8.1ms", "3.6ms", "2.3\u00d7"],
    ]

    y = Inches(1.2)
    col_specs = [(Inches(0.8), Inches(3.0)), (Inches(3.8), Inches(1.5)),
                 (Inches(5.3), Inches(1.5)), (Inches(6.8), Inches(1.5))]
    for col, (x, w) in enumerate(col_specs):
        add_text_box(slide, x, y, w, Inches(0.4),
                     headers[col], font_size=16, color=ACCENT, bold=True)

    for i, row in enumerate(rows):
        y_row = Inches(1.7 + i * 0.5)
        for col, (x, w) in enumerate(col_specs):
            c = ORANGE if col == 0 else (ACCENT2 if col == 3 else WHITE)
            add_text_box(slide, x, y_row, w, Inches(0.4),
                         row[col], font_size=15, color=c)

    # Key techniques
    add_text_box(slide, Inches(0.6), Inches(4.0), Inches(9.0), Inches(0.4),
                 "Key Optimization Techniques", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(4.5), Inches(8.5), Inches(2.5), [
        "Replaced spsolve with solve_banded for tridiagonal systems",
        "Vectorized matrix assembly (P2 FEM, Chebyshev)",
        "Precomputed geometric factors outside time loops",
        "Broadcasting operations instead of Python loops",
    ], font_size=16)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Project Architecture", font_size=32, color=ACCENT, bold=True)

    modules = [
        ("app/", "CLI entrypoint, benchmark runner"),
        ("solvers/", "FDM, FEM, FVM, Spectral, PINN (8 methods)"),
        ("features/", "Gradient, Laplacian, chi, energy, ML features"),
        ("metrics/", "L2 and L\u221e error computation"),
        ("policy/", "Post-hoc selection + ML decision tree + training"),
        ("reports/", "CSV and Markdown output generation"),
        ("tests/", "22 unit tests (pytest)"),
        ("docs/", "Manual, tutorial, slides, figures"),
    ]

    y = Inches(1.2)
    for mod, desc in modules:
        add_text_box(slide, Inches(0.8), y, Inches(2.0), Inches(0.4),
                     mod, font_size=16, color=ORANGE, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(2.9), y, Inches(6.8), Inches(0.4),
                     desc, font_size=15, color=WHITE)
        y += Inches(0.48)

    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(8.4), Inches(0.8),
                 "Dependencies: numpy, scipy  |  Optional: torch (PINN)\n"
                 "CI: GitHub Actions (Python 3.10 / 3.11 / 3.12)",
                 font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_cli_reference(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "CLI Reference", font_size=32, color=ACCENT, bold=True)

    commands = [
        ("make test", "Run all 22 tests"),
        ("make benchmark", "Run benchmark (default \u03b1)"),
        ("make train", "Generate data + train ML model"),
        ("--alpha 0.5 1.0", "Custom \u03b1 values"),
        ("--nr 101 --dt 0.0005", "Higher resolution"),
        ("--generate-data", "Create training CSV (~432 samples)"),
        ("--use-ml-selector", "Predict best solver, run only that"),
        ("--update", "Append results + retrain model"),
        ("--model-path PATH", "Custom model file path"),
    ]

    y = Inches(1.1)
    for cmd, desc in commands:
        add_text_box(slide, Inches(0.6), y, Inches(4.0), Inches(0.35),
                     cmd, font_size=14, color=ACCENT2, font_name="Consolas")
        add_text_box(slide, Inches(4.8), y, Inches(5.0), Inches(0.35),
                     desc, font_size=14, color=WHITE)
        y += Inches(0.42)


def slide_dev_process_1(prs):
    """Development process: Steps 1-12 (foundation)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Development Process (1/3): Foundation", font_size=32, color=ACCENT, bold=True)

    steps = [
        ("Steps 1-2", "Scaffold & Features",
         "pyproject.toml, Makefile, feature extraction (gradient, laplacian, chi, energy)"),
        ("Steps 3-6", "Three Solvers",
         "Implicit FDM (Crank-Nicolson), Spectral (cosine basis), PINN stub"),
        ("Steps 7-10", "Metrics & Policy",
         "L2/L\u221e error, solver selection (score = error + \u03bb\u00d7time), CSV/Markdown reports"),
        ("Steps 11-12", "Documentation & Run",
         "README, CLAUDE.md, first benchmark run \u2192 implicit_fdm wins for all \u03b1"),
    ]

    y = Inches(1.2)
    for label, title, desc in steps:
        add_text_box(slide, Inches(0.6), y, Inches(1.8), Inches(0.4),
                     label, font_size=14, color=ACCENT2, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(2.5), y, Inches(2.5), Inches(0.4),
                     title, font_size=16, color=ORANGE, bold=True)
        add_text_box(slide, Inches(5.0), y, Inches(4.8), Inches(0.5),
                     desc, font_size=14, color=WHITE)
        y += Inches(0.7)

    add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(8.8), Inches(0.5),
                 "Key Design Decisions", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), y + Inches(0.85), Inches(8.4), Inches(2.5), [
        "\u2022  r=0 singularity: L'H\u00f4pital rule \u2192 (1/r)\u2202(r\u03c7\u2202T/\u2202r) \u2192 2\u03c7\u2202\u00b2T/\u2202r\u00b2",
        "\u2022  Reference solution: same FDM at 4\u00d7 refinement (no analytical solution needed)",
        "\u2022  Modular SolverBase ABC \u2192 easy to add new solvers",
        "\u2022  Minimal dependencies: numpy-only core",
    ], font_size=15)


def slide_dev_process_2(prs):
    """Development process: Steps 13-17 (optimization & tooling)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Development Process (2/3): Optimization", font_size=32, color=ACCENT, bold=True)

    steps = [
        ("Step 13", "Detailed Manual",
         "docs/MANUAL.md \u2014 reference solution process, error definitions, PDE formulation"),
        ("Step 14", "Claude Code Skills",
         "5 slash commands (/run-benchmark, /add-solver, etc.) + auto syntax-check hook"),
        ("Step 15", "English Translation",
         "All docs and skill files translated from Japanese to English"),
        ("Step 16", "Vectorization",
         "FDM ~2.6\u00d7 speedup, Spectral ~23\u00d7 speedup via matrix ops"),
        ("Step 17", "scipy.linalg.solve_banded",
         "Replaced Python-loop Thomas algorithm with LAPACK \u2192 additional ~2.5\u00d7 speedup"),
    ]

    y = Inches(1.2)
    for label, title, desc in steps:
        add_text_box(slide, Inches(0.6), y, Inches(1.5), Inches(0.4),
                     label, font_size=14, color=ACCENT2, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(2.2), y, Inches(2.8), Inches(0.4),
                     title, font_size=16, color=ORANGE, bold=True)
        add_text_box(slide, Inches(5.0), y, Inches(4.8), Inches(0.5),
                     desc, font_size=14, color=WHITE)
        y += Inches(0.65)

    add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(8.8), Inches(0.5),
                 "Performance Gains", font_size=20, color=ACCENT2, bold=True)

    # Performance comparison
    add_text_box(slide, Inches(0.8), y + Inches(0.85), Inches(8.4), Inches(0.4),
                 "Implicit FDM total speedup: ~6.5\u00d7   |   Spectral total speedup: ~23\u00d7",
                 font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.8), y + Inches(1.4), Inches(8.4), Inches(1.5), [
        "\u2022  Vectorized tridiagonal construction + precomputed geometric factors",
        "\u2022  Matrix-vector transforms replacing per-mode loops in spectral solver",
        "\u2022  LAPACK dgbsv (Fortran) replacing Python-loop Thomas algorithm",
    ], font_size=15)


def slide_dev_process_3(prs):
    """Development process: Step 18+ (ML & documentation)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Development Process (3/3): ML & Docs", font_size=32, color=ACCENT, bold=True)

    steps = [
        ("Step 18", "ML Solver Selector",
         "12 features, numpy-only decision tree, --generate-data / --use-ml-selector"),
        ("  +update", "Incremental Learning",
         "--update flag: append results to CSV + retrain \u2192 model improves over time"),
        ("Docs", "Tutorial & Figures",
         "6-part tutorial, 5 auto-generated figures (generate_figures.py)"),
        ("CI", "GitHub Actions",
         "Test matrix: Python 3.10 / 3.11 / 3.12 on every push and PR"),
        ("Fix", "Python 3.12 Compat",
         "build-backend: setuptools.backends._legacy \u2192 setuptools.build_meta"),
        ("Slides", "Presentation",
         "python-pptx auto-generated slides with embedded benchmark figures"),
    ]

    y = Inches(1.2)
    for label, title, desc in steps:
        add_text_box(slide, Inches(0.6), y, Inches(1.5), Inches(0.4),
                     label, font_size=14, color=ACCENT2, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(2.2), y, Inches(2.8), Inches(0.4),
                     title, font_size=16, color=ORANGE, bold=True)
        add_text_box(slide, Inches(5.0), y, Inches(4.8), Inches(0.5),
                     desc, font_size=14, color=WHITE)
        y += Inches(0.6)

    add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(8.8), Inches(0.5),
                 "Key Innovation: ML Selector", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), y + Inches(0.85), Inches(8.4), Inches(1.8), [
        "\u2022  No sklearn \u2014 CART tree with Gini impurity in pure numpy (~100 LOC)",
        "\u2022  Features designed for pre-solve prediction (problem params + T\u2080 properties)",
        "\u2022  Incremental learning closes the loop: use \u2192 observe \u2192 improve",
        "\u2022  All artifacts reproducible: python docs/generate_figures.py, python docs/generate_slides.py",
    ], font_size=15)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Summary", font_size=32, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.5), [
        "\u2705  Complete PDE benchmark suite for fusion heat transport",
        "\u2705  Three solvers with automated error evaluation",
        "\u2705  ML solver selector: numpy-only decision tree (14 features)",
        "\u2705  Incremental learning via --update flag",
        "\u2705  Comprehensive documentation: manual, tutorial, figures",
        "\u2705  CI/CD with GitHub Actions (Python 3.10-3.12)",
    ], font_size=20)

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.5),
                 "Future Directions", font_size=22, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(8.4), Inches(2.0), [
        "\u2022  Full PINN implementation with PyTorch",
        "\u2022  Additional solvers (explicit FDM, FEM, etc.)",
        "\u2022  2D extension for toroidal geometry",
        "\u2022  Ensemble model or gradient boosting for selector",
    ], font_size=16, color=LIGHT_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title_page(prs)
    slide_overview(prs)
    slide_pde(prs)
    slide_initial_conditions(prs)
    slide_solvers(prs)
    slide_benchmark_results(prs)
    slide_selection_policy(prs)
    slide_ml_features(prs)
    slide_ml_workflow(prs)
    slide_incremental(prs)
    slide_architecture(prs)
    slide_cli_reference(prs)
    slide_dev_process_1(prs)
    slide_dev_process_2(prs)
    slide_dev_process_3(prs)
    slide_summary(prs)

    prs.save(OUTPATH)
    print(f"Saved {len(prs.slides)} slides to {OUTPATH}")


if __name__ == "__main__":
    main()
