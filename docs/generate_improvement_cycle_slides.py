"""Generate presentation slides for the Multi-Agent Improvement Cycle."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPATH = os.path.join(os.path.dirname(__file__), "improvement_cycle_slides.pptx")

# Color palette (same as generate_slides.py)
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_MID = RGBColor(0x24, 0x24, 0x3E)
ACCENT = RGBColor(0x4E, 0xA8, 0xDE)
ACCENT2 = RGBColor(0x50, 0xC8, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
RED = RGBColor(0xFF, 0x60, 0x60)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_code_box(slide, left, top, width, height, code, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x28)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xA0, 0xE0, 0xA0)
    p.font.name = "Consolas"
    return tf


def add_box(slide, left, top, width, height, fill_color=BG_MID,
            border_color=ACCENT, border_width=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    return shape


def add_table_rows(slide, x0, y0, headers, rows, col_widths,
                   header_color=ACCENT, row_colors=None, font_size=13):
    """Add a table-like layout of text boxes."""
    # Header
    x = x0
    for i, h in enumerate(headers):
        add_text_box(slide, x, y0, col_widths[i], Inches(0.35),
                     h, font_size=font_size, color=header_color, bold=True)
        x += col_widths[i]

    # Rows
    for ri, row in enumerate(rows):
        y = y0 + Inches(0.38 * (ri + 1))
        x = x0
        for ci, cell in enumerate(row):
            if row_colors and ci < len(row_colors):
                c = row_colors[ci]
            else:
                c = WHITE
            add_text_box(slide, x, y, col_widths[ci], Inches(0.35),
                         cell, font_size=font_size - 1, color=c)
            x += col_widths[ci]


# =========================================================================
# Slide functions
# =========================================================================

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
                 "Multi-Agent Improvement Cycle",
                 font_size=38, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(8.4), Inches(0.7),
                 "Fusion Heat Transport PDE Benchmark",
                 font_size=28, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.0),
                 "Sequential Agent Pipeline: Pareto Analysis, Bottleneck Detection,\n"
                 "Proposal Generation & Multi-Perspective Evaluation\n"
                 "with PHYSBO Bayesian Optimization",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """Problem and motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Problem & Motivation", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(4.3), Inches(0.4),
                 "Challenge", font_size=20, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(4.1), Inches(2.5), [
        "\u25b6  8 numerical solvers (FDM, FEM, FVM, Spectral, PINN)",
        "\u25b6  Each solver has different accuracy/speed/stability",
        "\u25b6  Performance varies across \u03b1, dt, nr, IC type",
        "\u25b6  Manual tuning is tedious and error-prone",
    ], font_size=15)

    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.4),
                 "Solution: Multi-Agent Automation", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(5.4), Inches(1.6), Inches(4.3), Inches(2.5), [
        "\u25b6  Automated Pareto analysis across parameters",
        "\u25b6  Systematic bottleneck detection (8 categories)",
        "\u25b6  AI-generated improvement proposals",
        "\u25b6  Multi-perspective evaluation (4 views)",
    ], font_size=15)

    add_text_box(slide, Inches(0.6), Inches(4.3), Inches(8.8), Inches(0.5),
                 "PDE: \u2202T/\u2202t = (1/r) \u2202/\u2202r (r \u03c7(|\u2202T/\u2202r|) \u2202T/\u2202r)",
                 font_size=20, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")
    add_text_box(slide, Inches(0.6), Inches(4.9), Inches(8.8), Inches(0.5),
                 "\u03c7(|T'|) = (|T'| \u2212 0.5)^\u03b1 + 0.1  (|T'| > 0.5),  else \u03c7 = 0.1",
                 font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")


def slide_03_architecture(prs):
    """Multi-agent architecture overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Agent Architecture (Sequential Pipeline)",
                 font_size=32, color=ACCENT, bold=True)

    # Agent boxes with flow
    agents = [
        (0.3, 1.3, 2.8, 1.4, "Pareto Analysis\nAgent", ACCENT,
         "Parameter sweep\nPHYSBO optimization\nPareto front computation"),
        (3.5, 1.3, 2.8, 1.4, "Bottleneck\nAnalysis Agent", ORANGE,
         "8 detection categories\nSeverity scoring\nAffected solver mapping"),
        (6.7, 1.3, 2.8, 1.4, "Proposal\nGeneration Agent", ACCENT2,
         "1-3 proposals/bottleneck\n3 proposal types\nImplementation sketches"),
        (0.3, 3.5, 2.8, 1.4, "Evaluation Agent\n(4 Perspectives)", PURPLE,
         "Accuracy(1.5) Speed(1.0)\nStability(1.2) Complexity(0.8)\nWeighted scoring"),
        (3.5, 3.5, 2.8, 1.4, "Review &\nImplementation", LIGHT_GRAY,
         "Auto/interactive approve\nscore \u2265 4.0 \u2192 approve\nImpl. sketch output"),
        (6.7, 3.5, 2.8, 1.4, "Report Agent", ACCENT,
         "Markdown report\nCycle archive\nHistory tracking"),
    ]

    for x, y, w, h, title, color, desc in agents:
        box = add_box(slide, Inches(x), Inches(y), Inches(w), Inches(h),
                      border_color=color)
        add_text_box(slide, Inches(x + 0.1), Inches(y + 0.05),
                     Inches(w - 0.2), Inches(0.5),
                     title, font_size=13, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(y + 0.55),
                     Inches(w - 0.2), Inches(0.8),
                     desc, font_size=10, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Arrows between boxes (row 1)
    for x in [3.1, 6.3]:
        add_text_box(slide, Inches(x), Inches(1.7), Inches(0.4), Inches(0.4),
                     "\u25b6", font_size=24, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Arrow down from Proposal to Evaluation
    add_text_box(slide, Inches(4.6), Inches(2.8), Inches(0.5), Inches(0.6),
                 "\u25bc", font_size=24, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Arrows between boxes (row 2)
    for x in [3.1, 6.3]:
        add_text_box(slide, Inches(x), Inches(3.9), Inches(0.4), Inches(0.4),
                     "\u25b6", font_size=24, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Hypothesis system (bottom)
    add_box(slide, Inches(0.3), Inches(5.3), Inches(4.5), Inches(1.0),
            border_color=PURPLE)
    add_text_box(slide, Inches(0.5), Inches(5.35), Inches(4.1), Inches(0.35),
                 "HypothesisTracker + ExperimentRunner",
                 font_size=13, color=PURPLE, bold=True)
    add_text_box(slide, Inches(0.5), Inches(5.75), Inches(4.1), Inches(0.5),
                 "Hypothesis registration \u2192 Experiment \u2192 Verification \u2192 Confidence update",
                 font_size=11, color=WHITE)

    add_box(slide, Inches(5.2), Inches(5.3), Inches(4.3), Inches(1.0),
            border_color=ORANGE)
    add_text_box(slide, Inches(5.4), Inches(5.35), Inches(3.9), Inches(0.35),
                 "CycleCoordinator (Sequential Orchestrator)",
                 font_size=13, color=ORANGE, bold=True)
    add_text_box(slide, Inches(5.4), Inches(5.75), Inches(3.9), Inches(0.5),
                 "for phase in PHASES: run_phase() \u2022 State persistence \u2022 Resume",
                 font_size=11, color=WHITE)


def slide_04_phase1_pareto(prs):
    """Phase 1: Pareto analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Phase 1: Pareto Analysis", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(4.3), Inches(0.4),
                 "Per-Solver Analysis", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(4.1), Inches(2.0), [
        "\u25b6  Parameter sweep: \u03b1=[0.0, 0.5, 1.0]",
        "\u25b6  dt=[0.001, 0.0005], nr=61 (fixed)",
        "\u25b6  PHYSBO or grid-based dt exploration",
        "\u25b6  Stability check + L2/L\u221e error computation",
        "\u25b6  Pareto rank assignment (0=optimal)",
    ], font_size=15)

    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.4),
                 "Cross-Solver Comparison", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(5.4), Inches(1.6), Inches(4.3), Inches(2.0), [
        "\u25b6  Rankings per problem setting",
        "\u25b6  Win counts (accuracy, speed, Pareto)",
        "\u25b6  Coverage gap detection",
        "\u25b6  Overall solver rankings",
    ], font_size=15)

    # PHYSBO integration
    add_text_box(slide, Inches(0.6), Inches(3.6), Inches(8.8), Inches(0.4),
                 "PHYSBO Bayesian Optimization Integration", font_size=20,
                 color=ORANGE, bold=True)

    add_code_box(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(2.5),
                 "# Feature: log10(dt) (1-dim, 80 discrete candidates)\n"
                 "# Objectives: -L2_error, -wall_time (2-objective maximization)\n"
                 "# Discrete multi-objective optimization per (alpha, ic_type):\n\n"
                 "dt_candidates = np.logspace(-5, -2, 80)  # 80 log-spaced\n"
                 "test_X = np.log10(dt_candidates).reshape(-1, 1)\n"
                 "policy = physbo.search.discrete_multi.Policy(\n"
                 "    test_X=test_X, num_objectives=2)  # discrete search\n"
                 "policy.random_search(max_num_probes=5)   # random phase\n"
                 "policy.bayes_search(max_num_probes=15, score='HVPI')  # Bayesian",
                 font_size=11)


def slide_05_phase2_bottleneck(prs):
    """Phase 2: Bottleneck detection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Phase 2: Bottleneck Detection", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.4),
                 "8 Detection Categories", font_size=20, color=ACCENT2, bold=True)

    categories = [
        ("stability", "LOW", "Solver stability rate < 90%"),
        ("accuracy_gap", "MED", "Large L2 error differences between solvers"),
        ("speed_gap", "LOW", "Execution time ratio > 100\u00d7"),
        ("coverage_gap", "MED", "Single solver dominates Pareto front (>80%)"),
        ("no_stable_solver", "HIGH", "No solver stable for a problem setting"),
        ("solver_dominance", "MED", "One solver wins >80% of problems"),
        ("cross_accuracy_gap", "HIGH", "Best solver still has L2 > 0.5"),
        ("solver_instability", "MED", "Solver fails on >20% of problems"),
    ]

    y = Inches(1.7)
    for cat, sev, desc in categories:
        sev_color = RED if sev == "HIGH" else (ORANGE if sev == "MED" else ACCENT2)
        add_text_box(slide, Inches(0.6), y, Inches(2.4), Inches(0.35),
                     cat, font_size=13, color=ACCENT, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(3.0), y, Inches(0.7), Inches(0.35),
                     sev, font_size=12, color=sev_color, bold=True)
        add_text_box(slide, Inches(3.8), y, Inches(5.8), Inches(0.35),
                     desc, font_size=13, color=WHITE)
        y += Inches(0.4)

    add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(8.8), Inches(0.5),
                 "Each bottleneck includes: severity, affected solvers, "
                 "evidence dict, and suggested actions",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_06_phase34_proposal_eval(prs):
    """Phase 3-4: Proposal generation and evaluation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Phase 3-4: Proposal & Multi-Perspective Evaluation",
                 font_size=30, color=ACCENT, bold=True)

    # Left: Proposal Generation
    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(4.5), Inches(0.4),
                 "Phase 3: Proposal Generation", font_size=18, color=ACCENT2, bold=True)

    proposals = [
        ("parameter_tuning", "dt constraints, nr adjustment"),
        ("algorithm_tweak", "Adaptive time-stepping, lagged coefficients"),
        ("new_solver", "IMEX for stiff problems, hybrid methods"),
    ]
    y = Inches(1.7)
    for ptype, desc in proposals:
        add_text_box(slide, Inches(0.6), y, Inches(2.0), Inches(0.35),
                     ptype, font_size=13, color=ORANGE, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(2.7), y, Inches(2.2), Inches(0.35),
                     desc, font_size=12, color=WHITE)
        y += Inches(0.4)

    # Right: Evaluation Agent
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.4),
                 "Phase 4: 4-Perspective Evaluation", font_size=18,
                 color=ACCENT2, bold=True)

    perspectives = [
        ("Accuracy", "1.5", "Resolution, adaptive methods"),
        ("Speed", "1.0", "Optimization, vectorization"),
        ("Stability", "1.2", "Stability-focused proposals"),
        ("Complexity", "0.8", "Simpler = better (param > algo > new)"),
    ]
    y = Inches(1.7)
    for name, weight, desc in perspectives:
        add_text_box(slide, Inches(5.4), y, Inches(1.5), Inches(0.35),
                     f"{name} ({weight})", font_size=13, color=PURPLE, bold=True)
        add_text_box(slide, Inches(7.0), y, Inches(2.7), Inches(0.35),
                     desc, font_size=12, color=WHITE)
        y += Inches(0.4)

    # Note about sequential execution
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(8.8), Inches(0.35),
                 "4 Perspectives run sequentially (for loop). "
                 "Parallelizable via concurrent.futures if extended to heavy tasks.",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Scoring formula
    add_text_box(slide, Inches(0.6), Inches(3.9), Inches(8.8), Inches(0.5),
                 "overall = \u03a3(weight_i \u00d7 score_i) / \u03a3(weight_i)    "
                 "approve \u2265 4.0  |  consider \u2265 3.0  |  reject < 3.0",
                 font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")

    # Example evaluation
    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.4),
                 "Example: Multi-Agent Evaluation Scores (Cycle 1)",
                 font_size=16, color=ACCENT2, bold=True)

    headers = ["Proposal", "Accuracy", "Speed", "Stability", "Complex.", "Overall", "Rec."]
    rows = [
        ["P001: Adaptive stepping", "4.0", "3.0", "4.5", "2.5", "3.38", "consider"],
        ["P002: dt constraints", "3.5", "4.5", "4.0", "4.5", "3.81", "consider"],
        ["P003: Vectorize FVM", "3.0", "5.0", "3.0", "3.5", "3.38", "consider"],
    ]
    col_w = [Inches(2.5), Inches(0.9), Inches(0.8), Inches(0.9),
             Inches(0.9), Inches(0.9), Inches(1.0)]
    add_table_rows(slide, Inches(0.5), Inches(5.0), headers, rows, col_w,
                   row_colors=[ORANGE, WHITE, WHITE, WHITE, WHITE, ACCENT2, LIGHT_GRAY],
                   font_size=12)


def slide_07_phase567(prs):
    """Phase 5-7: Review, Implementation, Report."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Phase 5-7: Review \u2192 Implementation \u2192 Report",
                 font_size=30, color=ACCENT, bold=True)

    # Phase 5: Review
    add_box(slide, Inches(0.3), Inches(1.2), Inches(2.9), Inches(2.5),
            border_color=ACCENT)
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(2.7), Inches(0.35),
                 "Phase 5: Review", font_size=18, color=ACCENT, bold=True)
    add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(2.6), Inches(1.5), [
        "Auto: score \u2265 4.0 \u2192 approve",
        "Interactive: Y/n/q per proposal",
        "Status \u2192 approved / rejected",
    ], font_size=13)

    # Phase 6: Implementation
    add_box(slide, Inches(3.5), Inches(1.2), Inches(2.9), Inches(2.5),
            border_color=ACCENT2)
    add_text_box(slide, Inches(3.6), Inches(1.3), Inches(2.7), Inches(0.35),
                 "Phase 6: Implement", font_size=18, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(3.7), Inches(1.8), Inches(2.6), Inches(1.5), [
        "param_tuning: preview sketch",
        "algo_tweak: manual guidance",
        "new_solver: code template",
    ], font_size=13)

    # Phase 7: Report
    add_box(slide, Inches(6.7), Inches(1.2), Inches(2.9), Inches(2.5),
            border_color=ORANGE)
    add_text_box(slide, Inches(6.8), Inches(1.3), Inches(2.7), Inches(0.35),
                 "Phase 7: Report", font_size=18, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(6.9), Inches(1.8), Inches(2.6), Inches(1.5), [
        "Markdown cycle report",
        "cycle_NNN_datetime.md",
        "History JSON updated",
    ], font_size=13)

    # Arrows
    for x in [3.15, 6.35]:
        add_text_box(slide, Inches(x), Inches(2.2), Inches(0.4), Inches(0.4),
                     "\u25b6", font_size=24, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Report contents
    add_text_box(slide, Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.4),
                 "Report Contents", font_size=18, color=ACCENT2, bold=True)

    sections = [
        "Executive summary (solver count, bottleneck count, proposals, approved)",
        "Cross-solver analysis: rankings, win counts, per-problem results",
        "Per-solver Pareto analysis: stability rates, error/time ranges",
        "Identified bottlenecks with severity and suggested actions",
        "Proposals with rationale and multi-perspective evaluation scores",
        "Next cycle recommendations",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.5), Inches(8.4), Inches(2.5),
                    [f"\u25b6  {s}" for s in sections], font_size=14)


def slide_08_physbo(prs):
    """PHYSBO Bayesian optimization details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "PHYSBO: Bayesian Multi-Objective Optimization",
                 font_size=30, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.5),
                 "Feature: log10(dt) [80 discrete candidates]    "
                 "Objectives: -L2_error, -wall_time",
                 font_size=17, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")

    # Grid vs PHYSBO comparison
    headers = ["", "Grid Search", "PHYSBO"]
    rows = [
        ["Feature", "\u2014", "log10(dt) (1-dim)"],
        ["Objectives", "L2, wall_time", "-L2, -wall_time (2-obj max)"],
        ["Evaluations", "All dt \u00d7 all problems", "20/problem (5+15)"],
        ["Search space", "Fixed grid", "80 log-spaced candidates"],
        ["Strategy", "Exhaustive", "HVPI-guided exploration"],
        ["Requirement", "Always available", "physbo package needed"],
    ]
    col_w = [Inches(1.8), Inches(3.2), Inches(3.5)]
    add_table_rows(slide, Inches(0.6), Inches(1.8), headers, rows, col_w,
                   row_colors=[ACCENT2, WHITE, WHITE],
                   font_size=14)

    # Flow
    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.5),
                 "Per (alpha, ic_type) pair:", font_size=16, color=ACCENT2, bold=True)
    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.5),
                 "80 dt candidates  \u2192  5 random probes  \u2192  "
                 "15 HVPI-guided probes  \u2192  best stable point",
                 font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(5.7), Inches(8.8), Inches(0.5),
                 "Fallback: auto-detect physbo availability; "
                 "grid search if not installed",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_09_rankings(prs):
    """Results: Cross-solver rankings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Results: Cross-Solver Rankings", font_size=32, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.4),
                 "Tutorial Run Output (\u03b1 = 0.0, 0.5, 1.0)", font_size=16,
                 color=LIGHT_GRAY)

    headers = ["Rank", "Solver", "Avg Rank", "Stability", "Min L2 Error"]
    rows = [
        ["#1", "implicit_fdm", "1.3", "100%", "0.054"],
        ["#2", "cell_centered_fvm", "2.0", "100%", "0.058"],
        ["#3", "compact4_fdm", "2.7", "100%", "0.115"],
        ["#4", "p2_fem", "3.5", "100%", "0.108"],
        ["#5", "imex_fdm", "4.8", "100%", "0.492"],
        ["#6", "cosine_spectral", "5.2", "67%", "varies"],
        ["#7", "pinn_stub", "6.5", "\u2014", "\u2014"],
        ["#8", "chebyshev_spectral", "7.0", "\u2014", "\u2014"],
    ]
    col_w = [Inches(0.7), Inches(2.8), Inches(1.3), Inches(1.2), Inches(1.5)]
    add_table_rows(slide, Inches(0.8), Inches(1.5), headers, rows, col_w,
                   row_colors=[ACCENT2, ORANGE, WHITE, WHITE, WHITE],
                   font_size=14)

    add_text_box(slide, Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.5),
                 "implicit_fdm consistently ranks #1 across all \u03b1 values",
                 font_size=16, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), Inches(5.7), Inches(8.8), Inches(0.4),
                 "Reference: Implicit FDM with 4\u00d7 grid refinement (nr\u00d74, dt/4)",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_10_pareto_results(prs):
    """Results: Per-solver Pareto analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Results: Per-Solver Pareto Analysis", font_size=32,
                 color=ACCENT, bold=True)

    headers = ["Solver", "Total", "Stable", "Pareto-Opt", "Min Error", "Max Error"]
    rows = [
        ["implicit_fdm", "12", "12", "5", "0.054", "0.296"],
        ["cell_centered_fvm", "12", "12", "5", "0.058", "0.287"],
        ["compact4_fdm", "12", "12", "4", "0.115", "0.489"],
        ["p2_fem", "12", "12", "5", "0.045", "0.531"],
        ["imex_fdm", "12", "12", "3", "0.492", "1.203"],
        ["cosine_spectral", "12", "8", "3", "0.089", "0.812"],
        ["chebyshev_spectral", "12", "4", "2", "0.799", "1.532"],
        ["pinn_stub", "12", "0", "0", "\u2014", "\u2014"],
    ]
    col_w = [Inches(2.4), Inches(0.8), Inches(0.9), Inches(1.2),
             Inches(1.3), Inches(1.3)]
    add_table_rows(slide, Inches(0.6), Inches(1.2), headers, rows, col_w,
                   row_colors=[ORANGE, WHITE, WHITE, ACCENT2, WHITE, WHITE],
                   font_size=13)

    add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(1.5), [
        "\u25b6  FDM/FVM solvers: 100% stability, competitive error",
        "\u25b6  cosine_spectral: dt-sensitive stability (67%), best at low \u03b1",
        "\u25b6  p2_fem: good accuracy but 10-100\u00d7 slower",
    ], font_size=15)


def slide_11_bottleneck_results(prs):
    """Results: Detected bottlenecks and proposals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Results: Bottlenecks & Proposals", font_size=32,
                 color=ACCENT, bold=True)

    # Bottlenecks
    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.4),
                 "Detected Bottlenecks (Cycle 1)", font_size=18,
                 color=ORANGE, bold=True)

    bottlenecks = [
        ("speed_gap", "LOW",
         "fastest=2.98ms (cell_centered_fvm), slowest=2361.36ms \u2192 792\u00d7 gap"),
        ("stability", "MED",
         "cosine_spectral: 67% stability (\u03b1 \u2265 0.5 instability)"),
    ]
    y = Inches(1.6)
    for cat, sev, desc in bottlenecks:
        sev_color = RED if sev == "HIGH" else (ORANGE if sev == "MED" else ACCENT2)
        add_text_box(slide, Inches(0.6), y, Inches(1.8), Inches(0.35),
                     cat, font_size=13, color=ACCENT, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(2.5), y, Inches(0.6), Inches(0.35),
                     sev, font_size=12, color=sev_color, bold=True)
        add_text_box(slide, Inches(3.2), y, Inches(6.4), Inches(0.35),
                     desc, font_size=13, color=WHITE)
        y += Inches(0.45)

    # Proposals
    add_text_box(slide, Inches(0.6), Inches(2.8), Inches(8.8), Inches(0.4),
                 "Generated Proposals", font_size=18, color=ACCENT2, bold=True)

    headers = ["ID", "Type", "Title", "Score", "Recommendation"]
    rows = [
        ["P001", "algorithm_tweak", "Adaptive time-stepping for cosine_spectral",
         "3.38", "consider"],
        ["P002", "parameter_tuning", "Constrain dt for high-alpha problems",
         "3.81", "consider"],
        ["P003", "parameter_tuning", "Optimize FVM vectorization",
         "3.38", "consider"],
    ]
    col_w = [Inches(0.6), Inches(1.7), Inches(3.8), Inches(0.8), Inches(1.5)]
    add_table_rows(slide, Inches(0.5), Inches(3.3), headers, rows, col_w,
                   row_colors=[ACCENT2, ORANGE, WHITE, ACCENT2, LIGHT_GRAY],
                   font_size=12)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.5),
                 "All proposals scored 3.0-3.9 (\"consider\" range) \u2014 "
                 "no auto-approved proposals in cycle 1",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_12_hypothesis(prs):
    """Hypothesis-driven workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Hypothesis-Driven Workflow", font_size=32, color=ACCENT, bold=True)

    # Tracker description
    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(4.3), Inches(0.4),
                 "HypothesisTracker", font_size=20, color=PURPLE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(4.1), Inches(2.0), [
        "\u25b6  Register hypotheses with unique IDs",
        "\u25b6  Track verification history per hypothesis",
        "\u25b6  Auto-update status & confidence (0-1)",
        "\u25b6  JSON persistence for cross-session use",
    ], font_size=15)

    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.4),
                 "ExperimentRunner", font_size=20, color=PURPLE, bold=True)
    add_bullet_list(slide, Inches(5.4), Inches(1.6), Inches(4.3), Inches(2.0), [
        "\u25b6  Execute solver experiments with configs",
        "\u25b6  Compute reference (4\u00d7 refined ImplicitFDM)",
        "\u25b6  Record L2/L\u221e error, time, stability",
        "\u25b6  Append results to CSV database",
    ], font_size=15)

    # Status flow
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(8.8), Inches(0.5),
                 "untested  \u2192  experiment  \u2192  verified  \u2192  "
                 "confirmed / rejected / inconclusive",
                 font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Predefined experiments
    add_text_box(slide, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.4),
                 "Predefined Experiments", font_size=18, color=ACCENT2, bold=True)

    exps = [
        ("stability_map", "Sweep \u03b1\u00d7dt space (8 alphas, 5 dts)"),
        ("ic_comparison", "Compare across initial conditions"),
        ("pinn_comparison", "PINN variants vs FDM/Spectral"),
        ("linear_regime", "Test in purely linear regime (|dT/dr|<0.5)"),
        ("fine_sweep", "Exhaustive sweep (9\u03b1 \u00d7 5nr \u00d7 3dt \u00d7 3t_end)"),
    ]
    y = Inches(4.7)
    for name, desc in exps:
        add_text_box(slide, Inches(0.8), y, Inches(2.2), Inches(0.3),
                     name, font_size=13, color=ORANGE, font_name="Consolas")
        add_text_box(slide, Inches(3.1), y, Inches(6.5), Inches(0.3),
                     desc, font_size=13, color=WHITE)
        y += Inches(0.35)


def slide_13_hypothesis_example(prs):
    """Hypothesis verification examples."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Hypothesis Verification: Real Examples",
                 font_size=30, color=ACCENT, bold=True)

    # H1: Confirmed
    add_box(slide, Inches(0.3), Inches(1.1), Inches(9.2), Inches(2.3),
            border_color=ACCENT2)
    add_text_box(slide, Inches(0.5), Inches(1.15), Inches(8.8), Inches(0.4),
                 "H1: \"Smaller dt improves spectral solver stability\"",
                 font_size=16, color=ACCENT2, bold=True)
    add_text_box(slide, Inches(7.5), Inches(1.15), Inches(1.8), Inches(0.4),
                 "CONFIRMED", font_size=14, color=ACCENT2, bold=True)

    headers = ["dt", "0.0001", "0.0002", "0.0005", "0.001", "0.002"]
    rows = [["Stability", "100%", "100%", "100%", "50%", "0%"]]
    col_w = [Inches(1.2)] + [Inches(1.3)] * 5
    add_table_rows(slide, Inches(0.6), Inches(1.7), headers, rows, col_w,
                   row_colors=[ACCENT2] + [WHITE] * 5, font_size=14)

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(8.8), Inches(0.4),
                 "Confidence: 1.0  |  6+ verification attempts across cycles",
                 font_size=13, color=LIGHT_GRAY)

    # Compact4 vs Implicit
    add_box(slide, Inches(0.3), Inches(3.6), Inches(9.2), Inches(2.3),
            border_color=ORANGE)
    add_text_box(slide, Inches(0.5), Inches(3.65), Inches(8.8), Inches(0.4),
                 "H_compact4_best: \"Compact4 FDM beats implicit FDM at high \u03b1\"",
                 font_size=16, color=ORANGE, bold=True)
    add_text_box(slide, Inches(7.5), Inches(3.65), Inches(1.8), Inches(0.4),
                 "CONFIRMED", font_size=14, color=ACCENT2, bold=True)

    add_bullet_list(slide, Inches(0.6), Inches(4.2), Inches(8.6), Inches(1.5), [
        "\u25b6  \u03b1=1.0: compact4_fdm L2=0.115 vs implicit_fdm L2=0.296  (compact4 wins)",
        "\u25b6  \u03b1=1.5: compact4_fdm L2=0.189 vs implicit_fdm L2=0.450  (compact4 wins)",
        "\u25b6  4th-order spatial accuracy advantage grows with nonlinearity",
    ], font_size=14)


def slide_14_cycle_flow(prs):
    """Full 7-phase improvement cycle flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Improvement Cycle: Sequential 7-Phase Pipeline",
                 font_size=32, color=ACCENT, bold=True)

    phases = [
        ("1", "Pareto Analysis", ACCENT,
         "Parameter sweep + PHYSBO \u2192 Pareto fronts"),
        ("2", "Bottleneck Detection", ORANGE,
         "8 categories, severity scoring"),
        ("3", "Proposal Generation", ACCENT2,
         "1-3 proposals per bottleneck"),
        ("4", "Multi-Perspective Eval", PURPLE,
         "4 views, weighted scoring"),
        ("5", "Review", LIGHT_GRAY,
         "Auto (score\u22654.0) or interactive"),
        ("6", "Implementation", ACCENT2,
         "Sketches + guidance output"),
        ("7", "Report & Archive", ORANGE,
         "Markdown + history JSON"),
    ]

    for i, (num, title, color, desc) in enumerate(phases):
        y = Inches(1.2 + i * 0.8)

        # Phase number circle
        add_text_box(slide, Inches(0.5), y, Inches(0.5), Inches(0.5),
                     num, font_size=22, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, Inches(1.2), y, Inches(3.0), Inches(0.4),
                     title, font_size=17, color=color, bold=True)

        # Description
        add_text_box(slide, Inches(4.3), y, Inches(5.3), Inches(0.4),
                     desc, font_size=15, color=WHITE)

        # Arrow to next (except last)
        if i < len(phases) - 1:
            add_text_box(slide, Inches(0.55), y + Inches(0.4),
                         Inches(0.4), Inches(0.35),
                         "\u25bc", font_size=14, color=LIGHT_GRAY,
                         alignment=PP_ALIGN.CENTER)

    # CLI
    add_code_box(slide, Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.5),
                 "python docs/analysis/method_improvement_cycle.py --cycles 3 --auto",
                 font_size=13)


def slide_15_scoring_mechanism(prs):
    """How evaluation scores are actually determined."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Scoring Mechanism: Keyword-Based Rule Engine",
                 font_size=30, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.4),
                 "Each Perspective: default score=3.0, then if/else keyword matching on "
                 "proposal.title and proposal.proposal_type",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # AccuracyPerspective
    add_text_box(slide, Inches(0.4), Inches(1.5), Inches(4.5), Inches(0.35),
                 "AccuracyPerspective (weight=1.5)", font_size=14,
                 color=ACCENT2, bold=True)
    acc_rules = [
        'title \u2208 "resolution"|"accuracy" \u2192 4.5',
        'title \u2208 "adaptive"             \u2192 4.0',
        'type == parameter_tuning       \u2192 3.5',
        '(no match)                     \u2192 3.0',
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(4.3), Inches(1.5),
                    acc_rules, font_size=11, color=WHITE)

    # SpeedPerspective
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.35),
                 "SpeedPerspective (weight=1.0)", font_size=14,
                 color=ACCENT2, bold=True)
    spd_rules = [
        'title \u2208 "optimize"|"fast"      \u2192 4.5',
        'title \u2208 "adaptive"             \u2192 2.5',
        'title \u2208 "resolution"+"increase" \u2192 2.0',
        '(no match)                     \u2192 3.0',
    ]
    add_bullet_list(slide, Inches(5.4), Inches(1.9), Inches(4.3), Inches(1.5),
                    spd_rules, font_size=11, color=WHITE)

    # StabilityPerspective
    add_text_box(slide, Inches(0.4), Inches(3.4), Inches(4.5), Inches(0.35),
                 "StabilityPerspective (weight=1.2)", font_size=14,
                 color=ACCENT2, bold=True)
    stb_rules = [
        'title \u2208 "stability"|"adaptive" \u2192 5.0',
        'title \u2208 "constrain"            \u2192 4.5',
        'type == algorithm_tweak        \u2192 3.5',
        '(no match)                     \u2192 3.0',
    ]
    add_bullet_list(slide, Inches(0.6), Inches(3.8), Inches(4.3), Inches(1.5),
                    stb_rules, font_size=11, color=WHITE)

    # ComplexityPerspective
    add_text_box(slide, Inches(5.2), Inches(3.4), Inches(4.5), Inches(0.35),
                 "ComplexityPerspective (weight=0.8)", font_size=14,
                 color=ACCENT2, bold=True)
    cpx_rules = [
        'type == parameter_tuning  \u2192 4.5',
        'type == algorithm_tweak   \u2192 2.5',
        'type == new_solver        \u2192 1.5',
        'sketch > 20 lines         \u2192 score - 1.0',
    ]
    add_bullet_list(slide, Inches(5.4), Inches(3.8), Inches(4.3), Inches(1.5),
                    cpx_rules, font_size=11, color=WHITE)

    # Overall formula
    add_text_box(slide, Inches(0.6), Inches(5.3), Inches(8.8), Inches(0.5),
                 "overall = (acc\u00d71.5 + spd\u00d71.0 + stb\u00d71.2 + cpx\u00d70.8) / 4.5",
                 font_size=18, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")

    # Limitation note
    add_text_box(slide, Inches(0.6), Inches(5.9), Inches(8.8), Inches(0.8),
                 "Limitation: Scores are determined by keyword string matching, not by "
                 "understanding proposal content.\n"
                 "If no keyword matches, all 4 perspectives return 3.0 "
                 "\u2192 overall = 3.0 (fixed). "
                 "Score constants (e.g. 4.5, 2.5) are heuristic.",
                 font_size=13, color=RED)


def slide_16_advanced_multiagent(prs):
    """Introduction of advanced_multi_agent.py as future direction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Advanced Multi-Agent System (Prototype)",
                 font_size=30, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.4),
                 "docs/analysis/advanced_multi_agent.py \u2014 "
                 "Debate-based architecture with critical review",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Architecture diagram
    agents = [
        (0.3, 1.6, 2.1, 1.1, "Statistics\nAgent", ACCENT,
         "Distribution analysis\nDominance detection"),
        (2.7, 1.6, 2.1, 1.1, "Feature\nAgent", ACCENT,
         "Tree-based importance\nRedundancy detection"),
        (5.1, 1.6, 2.1, 1.1, "Hypothesis\nAgent", ACCENT2,
         "Alpha threshold test\nStiffness hypothesis\nGrid vs physics"),
        (7.5, 1.6, 2.1, 1.1, "Critic\nAgent", RED,
         "Challenges claims\nData diversity check"),
    ]

    for x, y, w, h, title, color, desc in agents:
        box = add_box(slide, Inches(x), Inches(y), Inches(w), Inches(h),
                      border_color=color)
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.05),
                     Inches(w - 0.1), Inches(0.4),
                     title, font_size=12, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.45),
                     Inches(w - 0.1), Inches(0.6),
                     desc, font_size=9, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Synthesis agent
    add_box(slide, Inches(2.5), Inches(3.0), Inches(5.0), Inches(0.8),
            border_color=PURPLE)
    add_text_box(slide, Inches(2.6), Inches(3.05), Inches(4.8), Inches(0.35),
                 "SynthesisAgent \u2014 Aggregates insights + hypotheses + critiques",
                 font_size=12, color=PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.6), Inches(3.4), Inches(4.8), Inches(0.3),
                 "Generates final report with key findings, confirmed hypotheses, caveats",
                 font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 4-phase flow
    add_text_box(slide, Inches(0.6), Inches(4.1), Inches(8.8), Inches(0.4),
                 "4-Phase Execution Flow", font_size=16,
                 color=ORANGE, bold=True)

    phases = [
        ("Phase 1", "Initial Analysis",
         "StatisticsAgent + FeatureAgent run independently on training data (X, y)"),
        ("Phase 2", "Hypothesis Testing",
         "HypothesisAgent tests 3 hypotheses: alpha threshold, stiffness, grid vs physics"),
        ("Phase 3", "Critical Review",
         "CriticAgent receives all insights and challenges weak claims"),
        ("Phase 4", "Synthesis",
         "SynthesisAgent merges findings into structured report"),
    ]
    y = Inches(4.6)
    for ph, title, desc in phases:
        add_text_box(slide, Inches(0.6), y, Inches(1.0), Inches(0.3),
                     ph, font_size=12, color=ACCENT2, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(1.7), y, Inches(2.0), Inches(0.3),
                     title, font_size=12, color=ORANGE, bold=True)
        add_text_box(slide, Inches(3.8), y, Inches(5.8), Inches(0.3),
                     desc, font_size=11, color=WHITE)
        y += Inches(0.35)

    # Key difference
    add_text_box(slide, Inches(0.6), Inches(6.1), Inches(8.8), Inches(0.8),
                 "Key difference from improvement cycle: CriticAgent can challenge "
                 "other agents' conclusions.\n"
                 "Still rule-based (no LLM), but demonstrates the debate pattern \u2014 "
                 "a step toward true multi-agent collaboration.",
                 font_size=13, color=LIGHT_GRAY)


def slide_17_summary(prs):
    """Summary and future directions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 "Summary & Future Directions", font_size=32,
                 color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.4),
                 "Key Features", font_size=20, color=ACCENT2, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(2.5), [
        "\u2705  6 specialized agents in sequential pipeline (clear I/O interfaces)",
        "\u2705  7-phase closed-loop improvement cycle (data-dependent, serial execution)",
        "\u2705  PHYSBO Bayesian discrete multi-objective optimization (dt search)",
        "\u2705  8-category bottleneck detection with severity scoring",
        "\u2705  4-perspective weighted evaluation (parallelizable for future heavy tasks)",
        "\u2705  Hypothesis-driven experiment framework with persistence",
        "\u2705  Full state persistence (JSON) for resume/reproducibility",
    ], font_size=16)

    add_text_box(slide, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.4),
                 "Design Principles", font_size=20, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(4.7), Inches(8.4), Inches(1.0), [
        "\u25b6  Hierarchical: CycleCoordinator \u2192 Agents \u2192 Data layer",
        "\u25b6  Modular: Each agent independently testable and replaceable",
        "\u25b6  Extensible: New solvers, bottleneck categories, eval perspectives",
    ], font_size=16)

    add_text_box(slide, Inches(0.6), Inches(5.6), Inches(8.8), Inches(0.4),
                 "Next Steps", font_size=20, color=PURPLE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(8.4), Inches(1.0), [
        "\u25b6  Parallelize Phase 1 (solver \u00d7 config) via ProcessPoolExecutor/MPI",
        "\u25b6  Expand PHYSBO search to multi-dimensional (dt + nr + t_end)",
        "\u25b6  Heavy Perspective evaluation \u2192 concurrent.futures parallelization",
    ], font_size=16, color=LIGHT_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_architecture(prs)
    slide_04_phase1_pareto(prs)
    slide_05_phase2_bottleneck(prs)
    slide_06_phase34_proposal_eval(prs)
    slide_07_phase567(prs)
    slide_08_physbo(prs)
    slide_09_rankings(prs)
    slide_10_pareto_results(prs)
    slide_11_bottleneck_results(prs)
    slide_12_hypothesis(prs)
    slide_13_hypothesis_example(prs)
    slide_14_cycle_flow(prs)
    slide_15_scoring_mechanism(prs)
    slide_16_advanced_multiagent(prs)
    slide_17_summary(prs)

    prs.save(OUTPATH)
    print(f"Saved {len(prs.slides)} slides to {OUTPATH}")


if __name__ == "__main__":
    main()
